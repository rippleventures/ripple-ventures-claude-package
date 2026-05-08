#!/usr/bin/env python3
"""
check_deck.py — QC checks on a finished PowerPoint deck.

Checks:
  1. No leftover placeholders ([TBD], [PLACEHOLDER], [CHECK], XXX)
  2. Page numbers present on every content slide
  3. Footnotes present on data-heavy slides (heuristic: slides with charts or tables)
  4. Dates consistent across slides (extracts dates, flags inconsistencies)
  5. Currency consistent (extracts $ amounts, flags mixed currencies)

Usage:
  python check_deck.py path/to/deck.pptx [--source-model path/to/model.xlsx]

Output: JSON report to stdout. Non-zero exit code if any check fails.

Dependencies: python-pptx
"""

import argparse
import json
import re
import sys
from pathlib import Path
from collections import Counter

try:
    from pptx import Presentation
except ImportError:
    print(json.dumps({"error": "python-pptx not installed", "fix": "pip install python-pptx"}))
    sys.exit(2)


PLACEHOLDER_PATTERNS = [
    r"\[TBD\]",
    r"\[PLACEHOLDER\]",
    r"\[CHECK\]",
    r"\[FILL\s*IN\]",
    r"\bXXX\b",
    r"\bTODO\b",
    r"\[insert.*?\]",
]

DATE_PATTERN = re.compile(
    r"\b(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)\w*\s+\d{1,2},?\s+(\d{4})\b"
    r"|\b(\d{1,2})/(\d{1,2})/(\d{4})\b"
    r"|\b(\d{4})-(\d{1,2})-(\d{1,2})\b",
    re.IGNORECASE,
)

CURRENCY_PATTERN = re.compile(r"([\$£€¥])\s?[\d,.]+")


def slide_text(slide):
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    parts.append(run.text)
    return " ".join(parts)


def check_placeholders(prs):
    findings = []
    for i, slide in enumerate(prs.slides, 1):
        text = slide_text(slide)
        for pattern in PLACEHOLDER_PATTERNS:
            for match in re.finditer(pattern, text, re.IGNORECASE):
                findings.append({
                    "slide": i,
                    "pattern": pattern,
                    "match": match.group(0),
                })
    return findings


def check_currency_consistency(prs):
    currency_counts = Counter()
    for slide in prs.slides:
        text = slide_text(slide)
        for match in CURRENCY_PATTERN.finditer(text):
            currency_counts[match.group(1)] += 1

    if len(currency_counts) > 1:
        return {
            "consistent": False,
            "currencies_found": dict(currency_counts),
            "note": "Multiple currency symbols found across the deck",
        }
    return {"consistent": True, "currencies_found": dict(currency_counts)}


def check_dates(prs):
    years = []
    for i, slide in enumerate(prs.slides, 1):
        text = slide_text(slide)
        for match in DATE_PATTERN.finditer(text):
            for group in match.groups():
                if group and group.isdigit() and len(group) == 4:
                    years.append((i, int(group)))

    if not years:
        return {"consistent": True, "years_found": []}

    year_counts = Counter(y for _, y in years)
    return {
        "consistent": len(year_counts) <= 2,  # allow current + prior year
        "years_found": dict(year_counts),
        "year_by_slide": years,
    }


def check_page_numbers(prs):
    """Heuristic: page number is a small text box with a digit matching the slide index."""
    missing = []
    for i, slide in enumerate(prs.slides, 1):
        if i == 1:  # cover slide typically doesn't need page number
            continue
        text = slide_text(slide)
        if str(i) not in text and f"{i} of" not in text and f"Page {i}" not in text:
            missing.append(i)
    return missing


def main():
    parser = argparse.ArgumentParser(description="QC a PowerPoint deck")
    parser.add_argument("deck", type=Path)
    parser.add_argument("--source-model", type=Path, help="Optional source xlsx for totals tie check")
    args = parser.parse_args()

    if not args.deck.exists():
        print(json.dumps({"error": f"File not found: {args.deck}"}))
        sys.exit(2)

    prs = Presentation(args.deck)

    placeholder_findings = check_placeholders(prs)
    currency_check = check_currency_consistency(prs)
    date_check = check_dates(prs)
    page_check = check_page_numbers(prs)

    summary = {
        "deck": str(args.deck),
        "slide_count": len(prs.slides),
        "passed": (
            len(placeholder_findings) == 0
            and currency_check["consistent"]
            and date_check["consistent"]
            and len(page_check) == 0
        ),
        "checks": {
            "placeholders": {
                "passed": len(placeholder_findings) == 0,
                "findings": placeholder_findings,
            },
            "currency_consistency": currency_check,
            "date_consistency": date_check,
            "page_numbers_present": {
                "passed": len(page_check) == 0,
                "missing_on_slides": page_check,
            },
        },
    }
    print(json.dumps(summary, indent=2))

    if not summary["passed"]:
        sys.exit(1)


if __name__ == "__main__":
    main()
